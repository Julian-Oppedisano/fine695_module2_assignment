from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
import os
from portfolio_utils import calculate_performance_metrics # To calculate SPY metrics

SLIDES_DIR = 'slides'
TEMPLATE_PATH = os.path.join(SLIDES_DIR, 'Module2_template.pptx')
OUTPUT_DIR = 'outputs'
CUM_RET_PNG = os.path.join(OUTPUT_DIR, 'cum_returns.png')
BEST_ALG_PATH = os.path.join(OUTPUT_DIR, 'best_alg.txt')
METRICS_PATH = os.path.join(OUTPUT_DIR, 'metrics.csv')
PERF_SUMMARY_PATH = os.path.join(OUTPUT_DIR, 'perf_summary.csv')
SPY_DATA_PATH = 'Data/SPY returns.xlsx'

def get_best_model_name():
    try:
        with open(BEST_ALG_PATH, 'r') as f:
            return f.read().strip()
    except FileNotFoundError:
        print(f"Error: {BEST_ALG_PATH} not found. Run select_best_model.py first.")
        return None

def load_model_data(model_name, metrics_path, summary_path):
    data = {'name': model_name, 'oos_r2': np.nan, 'perf_summary': pd.Series(dtype='float64')}
    try:
        metrics_df = pd.read_csv(metrics_path)
        model_r2_series = metrics_df[metrics_df['model'] == model_name]['oos_r2']
        if not model_r2_series.empty:
            data['oos_r2'] = model_r2_series.iloc[0]
    except (FileNotFoundError, IndexError, KeyError) as e:
        print(f"Error loading R2 for {model_name} from {metrics_path}: {e}")

    try:
        summary_df = pd.read_csv(summary_path)
        model_perf_series = summary_df[summary_df['model'] == model_name]
        if not model_perf_series.empty:
            data['perf_summary'] = model_perf_series.iloc[0].copy()
            # Ensure avg_return and volatility are present, will be filled later if needed
            if 'avg_return' not in data['perf_summary']:
                 data['perf_summary']['avg_return'] = np.nan
            if 'volatility' not in data['perf_summary']:
                 data['perf_summary']['volatility'] = np.nan
            
            # Calculate Ann. Alpha if monthly alpha exists
            if pd.notna(data['perf_summary'].get('alpha')):
                data['perf_summary']['alpha_ann'] = data['perf_summary'].get('alpha') * 12
            else:
                data['perf_summary']['alpha_ann'] = np.nan

    except (FileNotFoundError, IndexError, KeyError) as e:
        print(f"Error loading performance summary for {model_name} from {summary_path}: {e}")
    return data


def get_spy_performance_for_period(model_returns_df, spy_data_path=SPY_DATA_PATH):
    try:
        spy_returns_raw = pd.read_excel(spy_data_path, index_col=0)
        spy_returns_raw.index = pd.to_datetime(spy_returns_raw.index)
        spy_returns_raw.columns = ['spy_ret'] 

        # Determine alignment period: if model_returns_df is provided, use its period
        # Otherwise, use the full SPY period for SPY's own metrics.
        if model_returns_df is not None and not model_returns_df.empty:
            model_start_date = model_returns_df.index.min()
            model_end_date = model_returns_df.index.max()
            spy_aligned = spy_returns_raw[(spy_returns_raw.index >= model_start_date) & (spy_returns_raw.index <= model_end_date)]
            if spy_aligned.empty:
                print("SPY returns have no overlap with model OOS period. Using full SPY history for SPY metrics.")
                spy_aligned = spy_returns_raw # Fallback to full history if no overlap
        else:
            print("No model returns for alignment. Calculating SPY metrics over its entire history.")
            spy_aligned = spy_returns_raw
        
        spy_perf_metrics_calc = calculate_performance_metrics(spy_aligned, spy_aligned, 'SPY')
        
        spy_metrics_output = pd.Series(dtype='object')
        spy_metrics_output['name'] = "SPY"
        spy_metrics_output['alpha'] = 0.0 # Alpha of SPY vs SPY is 0
        spy_metrics_output['alpha_ann'] = 0.0
        spy_metrics_output['sharpe'] = spy_perf_metrics_calc.get('sharpe', np.nan).iloc[0] if pd.notna(spy_perf_metrics_calc.get('sharpe', np.nan).iloc[0]) else np.nan
        spy_metrics_output['avg_return'] = spy_aligned['spy_ret'].mean()
        spy_metrics_output['volatility'] = spy_aligned['spy_ret'].std()
        spy_metrics_output['max_drawdown'] = spy_perf_metrics_calc.get('max_drawdown', np.nan).iloc[0] if pd.notna(spy_perf_metrics_calc.get('max_drawdown', np.nan).iloc[0]) else np.nan
        spy_metrics_output['max_loss'] = spy_perf_metrics_calc.get('max_loss', np.nan).iloc[0] if pd.notna(spy_perf_metrics_calc.get('max_loss', np.nan).iloc[0]) else np.nan
        return spy_metrics_output

    except Exception as e:
        print(f"Error calculating SPY performance: {e}")
        return pd.Series({'name': "SPY", 'alpha':0.0, 'alpha_ann':0.0, 'sharpe':np.nan, 'avg_return':np.nan, 'volatility':np.nan, 'max_drawdown':np.nan, 'max_loss':np.nan}, dtype='object')

def update_presentation():
    best_model_name = get_best_model_name()
    if not best_model_name:
        print("Best model name not found. Exiting presentation update.")
        return

    # Load data for best model, IPCA, and Autoencoder
    best_model_data = load_model_data(best_model_name, METRICS_PATH, PERF_SUMMARY_PATH)
    ipca_data = load_model_data('ipca', METRICS_PATH, PERF_SUMMARY_PATH)
    autoencoder_data = load_model_data('autoencoder', METRICS_PATH, PERF_SUMMARY_PATH)

    # Fill avg_return and volatility for each loaded model from their specific port_ret files
    models_to_fill = {best_model_name: best_model_data, 'ipca': ipca_data, 'autoencoder': autoencoder_data}
    for name, data in models_to_fill.items():
        try:
            port_ret_path = os.path.join(OUTPUT_DIR, f'{name}_overall_port_ret.csv')
            returns_df = pd.read_csv(port_ret_path, index_col=0, parse_dates=True)
            returns_col = returns_df.columns[0]
            data['perf_summary']['avg_return'] = returns_df[returns_col].mean()
            data['perf_summary']['volatility'] = returns_df[returns_col].std()
        except FileNotFoundError:
            print(f"Portfolio returns file not found for {name}: {port_ret_path}. Avg_return/volatility will be NaN.")
            # Ensure keys exist even if file not found, if perf_summary itself is valid
            if isinstance(data.get('perf_summary'), pd.Series):
                 if 'avg_return' not in data['perf_summary']:
                      data['perf_summary']['avg_return'] = np.nan
                 if 'volatility' not in data['perf_summary']:
                      data['perf_summary']['volatility'] = np.nan
        except Exception as e:
            print(f"Error processing portfolio returns for {name}: {e}")

    best_model_returns_df = None
    try:
        best_model_port_ret_path = os.path.join(OUTPUT_DIR, f'{best_model_name}_overall_port_ret.csv')
        best_model_returns_df = pd.read_csv(best_model_port_ret_path, index_col=0, parse_dates=True)
    except FileNotFoundError:
        print(f"Best model ({best_model_name}) returns file not found: {best_model_port_ret_path}")
        
    spy_perf_data = get_spy_performance_for_period(best_model_returns_df) # Align SPY to best model's period

    # Slide 1: Methodology
    methodology_text = (
        f"• Models Evaluated: IPCA, Autoencoder, Neural Network (NN2), and other baseline models (Lasso, Ridge, ElasticNet, Decision Tree).\n"
        f"• Best Performing Strategy: {best_model_name.upper()} selected based on overall portfolio metrics (e.g., Sharpe Ratio).\n"
        f"• Core Approach: Expanding window for training (initial 10yr), validation (initial 2yr), and OOS testing (2017-2023), with yearly model updates.\n"
        f"• Portfolio: Equal-weighted top 50 stocks (from predictions), monthly rebalancing.\n"
        f"• Features: 145 lagged stock characteristics to predict next-month excess returns.\n"
        f"• Benchmark: S&P 500 (SPY). Environment: Frictionless trading assumed as per assignment scope."
    )
    
    # Slide 2: Out-of-Sample R²
    r2_text_lines = [f"Out-of-Sample R² (Average over 2017-2023 OOS Period):\n"]
    r2_text_lines.append(f"• {best_model_data['name'].upper()}: {best_model_data['oos_r2']:.4f}")
    r2_text_lines.append(f"• IPCA: {ipca_data['oos_r2']:.4f}")
    r2_text_lines.append(f"• Autoencoder: {autoencoder_data['oos_r2']:.4f}")
    r2_text_lines.append("\nInterpretation & Diagnostics:")
    r2_text_lines.append(f"• The {best_model_name.upper()} model shows an average OOS R² of {best_model_data['oos_r2']:.4f}. "
                         "Negative R² values indicate underperformance against a simple mean forecast.")
    if best_model_name == 'nn2' and pd.notna(best_model_data['oos_r2']) and best_model_data['oos_r2'] < -10: # Specific diagnostic for NN2 if R2 is very poor
         r2_text_lines.append("• For NN2, this average was heavily impacted by extremely poor predictive performance in specific years (e.g., 2017), "
                              "masking more moderate performance in other periods. This suggests instability in the NN2 model's predictions.")
    r2_text_lines.append("• This highlights the challenge of achieving consistently positive R² in stock return prediction.")
    oos_r2_full_text = '\n'.join(r2_text_lines)

    # Slide 5: Summary
    summary_text = "Performance summary for the best model is incomplete."
    if isinstance(best_model_data.get('perf_summary'), pd.Series) and not best_model_data['perf_summary'].empty:
        avg_ret_val = best_model_data['perf_summary'].get('avg_return', np.nan)
        sharpe_val = best_model_data['perf_summary'].get('sharpe', np.nan)
        oos_r2_val = best_model_data['oos_r2']
        summary_text = (
            f"The {best_model_name.upper()} strategy, despite a challenging average OOS R² of {oos_r2_val:.2f}, "
            f"yielded an annualized Sharpe ratio of {sharpe_val:.2f} and average monthly returns of {avg_ret_val*100:.2f}%. "
            f"Improvements should target predictive model stability and explore alternative risk/portfolio construction techniques."
        )
        if len(summary_text) > 250: # Heuristic to keep it punchy
            summary_text = (
                f"The {best_model_name.upper()} strategy (Sharpe: {sharpe_val:.2f}, Avg Rtn: {avg_ret_val*100:.2f}%) showed portfolio value despite predictive R² challenges ({oos_r2_val:.2f}). "
                f"Focus future work on model stability & portfolio optimization."
            )


    prs = Presentation(TEMPLATE_PATH)
    
    # Slide 1: Methodology
    slide = prs.slides[0]
    title_shape = slide.shapes.title
    if title_shape: title_shape.text = "Methodology Overview"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame; tf.clear()
        p = tf.add_paragraph(); p.text = methodology_text.replace("\\n", "\n"); p.font.size = Pt(16)
    
    # Slide 2: OOS R² of Key Models
    slide = prs.slides[1]
    title_shape = slide.shapes.title
    if title_shape: title_shape.text = "Out-of-Sample R² Analysis (Key Models)"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame; tf.clear()
        # Use multiple paragraphs for better spacing control
        content_parts = oos_r2_full_text.split('\n')
        for part in content_parts:
            p = tf.add_paragraph()
            p.text = part
            p.font.size = Pt(16)
            if part.startswith("Out-of-Sample R²") : p.font.bold = True; p.font.size = Pt(18)
            if part.startswith("•") : p.level = 1; p.font.size = Pt(15)

    # Slide 3: Performance Table
    slide = prs.slides[2]
    title_shape = slide.shapes.title
    if title_shape: title_shape.text = "Portfolio Performance Statistics (OOS 2017-2023)"
    for shape in slide.shapes: # Clear old table
        if shape.has_table: slide.shapes._spTree.remove(shape._element)
            
    table_models = [best_model_data, ipca_data, autoencoder_data, spy_perf_data]
    header = ["Metric", best_model_name.upper(), "IPCA", "Autoencoder", "SPY", f"{best_model_name.upper()} vs SPY Δ"]
    
    rows_data = []
    metrics_to_display = [
        ("Alpha (monthly)", 'alpha', '.4f', False),
        ("Alpha (annualized)", 'alpha_ann', '.4f', False),
        ("Sharpe Ratio (ann.)", 'sharpe', '.2f', False),
        ("Avg Return (monthly)", 'avg_return', '.2%', True), # True for percentage
        ("Volatility (monthly)", 'volatility', '.2%', True),
        ("Max Drawdown", 'max_drawdown', '.2%', True),
        ("Max 1-mo Loss", 'max_loss', '.2%', True)
    ]

    for display_name, key, fmt, is_percent in metrics_to_display:
        row = [display_name]
        best_model_val = np.nan
        spy_val_for_delta = np.nan

        for i, model_data_item in enumerate(table_models):
            val = model_data_item.get('perf_summary', pd.Series(dtype='float64')).get(key, np.nan) if i < 3 else model_data_item.get(key, np.nan)
            if is_percent and pd.notna(val): formatted_val = f"{val:{fmt}}" 
            elif pd.notna(val): formatted_val = f"{val:{fmt}}"
            else: formatted_val = "N/A"
            row.append(formatted_val)
            if i == 0: best_model_val = val # Best model's value
            if model_data_item.get('name') == "SPY": spy_val_for_delta = val
        
        # Calculate Delta (Best Model - SPY)
        if pd.notna(best_model_val) and pd.notna(spy_val_for_delta) and key not in ['alpha', 'alpha_ann']: # Alpha for SPY is 0, so delta is just best_model_val
            delta = best_model_val - spy_val_for_delta
            if is_percent: row.append(f"{delta:{fmt}}")
            else: row.append(f"{delta:{fmt}}")
        elif key in ['alpha', 'alpha_ann'] and pd.notna(best_model_val):
             delta = best_model_val # Since SPY alpha is 0
             row.append(f"{delta:{fmt}}")
        else:
            row.append("N/A")
        rows_data.append(row)

    if rows_data:
        num_rows = len(rows_data) + 1 # +1 for header
        num_cols = len(header)
        left, top, width = Inches(0.2), Inches(1.2), Inches(9.6)
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(num_rows * 0.35))
        table = table_shape.table

        for c, header_text in enumerate(header):
            table.cell(0,c).text = header_text
            p = table.cell(0,c).text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER

        for r_idx, data_row in enumerate(rows_data):
            for c_idx, cell_val in enumerate(data_row):
                table.cell(r_idx+1, c_idx).text = str(cell_val)
                p = table.cell(r_idx+1, c_idx).text_frame.paragraphs[0]; p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER
        
        # Adjust column widths (example)
        widths = [Inches(1.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.6)]
        for i, col_width in enumerate(widths):
            if i < num_cols: table.columns[i].width = col_width
    else:
        tf = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1)).text_frame
        p = tf.add_paragraph(); p.text = "Performance table data not available."; p.alignment = PP_ALIGN.CENTER

    # Slide 4: Cumulative Returns Plot
    slide = prs.slides[3]
    title_shape = slide.shapes.title
    if title_shape: 
        title_shape.text = f"Cumulative OOS Returns: {best_model_name.upper()} vs SPY (2017-2023)"

    if os.path.exists(CUM_RET_PNG):
        pic_left, pic_top, pic_width = Inches(0.5), Inches(1.2), Inches(9.0)
        for shape_idx in reversed(range(len(slide.shapes))):
            shape = slide.shapes[shape_idx]
            if shape.shape_type == 13: slide.shapes._spTree.remove(shape._element)
        slide.shapes.add_picture(CUM_RET_PNG, pic_left, pic_top, width=pic_width)
    else:
        tf = slide.placeholders[1] if len(slide.placeholders) > 1 else slide.shapes.add_textbox(Inches(1),Inches(1),Inches(8),Inches(1)).text_frame
        tf.clear(); p = tf.add_paragraph(); p.text = "Cumulative returns plot not found."; p.alignment = PP_ALIGN.CENTER
    
    # Slide 5: Summary & Future Directions
    slide = prs.slides[4]
    title_shape = slide.shapes.title
    if title_shape: title_shape.text = "Key Insights & Future Directions"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame; tf.clear()
        p = tf.add_paragraph(); p.text = summary_text.replace("\\n", "\n"); p.font.size = Pt(16)
    
    output_filename = f'Module3_Report_{best_model_name}.pptx' # Original name for now
    # If a generic name is required: output_filename = 'Module3_Report.pptx'
    output_ppt_path = os.path.join(OUTPUT_DIR, output_filename)
    prs.save(output_ppt_path)
    print(f"Presentation updated and saved to: {output_ppt_path}")

if __name__ == "__main__":
    update_presentation() 