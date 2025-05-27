from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
import os
from portfolio_utils import calculate_performance_metrics # To calculate SPY metrics

SLIDES_DIR = 'slides'
TEMPLATE_PATH = os.path.join(SLIDES_DIR, 'Module2_template.pptx')
OUTPUT_DIR = 'outputs'
CUM_RET_PNG = os.path.join(OUTPUT_DIR, 'cum_returns.png')
BEST_ALG_PATH = os.path.join(OUTPUT_DIR, 'best_alg.txt')
METRICS_PATH = os.path.join(OUTPUT_DIR, 'metrics.csv')
PERF_SUMMARY_PATH = os.path.join(OUTPUT_DIR, 'perf_summary.csv')
SPY_DATA_PATH = 'Data/SPY returns.xlsx'

def get_best_model_name():
    try:
        with open(BEST_ALG_PATH, 'r') as f:
            return f.read().strip()
    except FileNotFoundError:
        print(f"Error: {BEST_ALG_PATH} not found. Run select_best_model.py first.")
        return None

def load_model_metrics(best_model_name):
    try:
        metrics_df = pd.read_csv(METRICS_PATH)
        model_r2 = metrics_df[metrics_df['model'] == best_model_name]['oos_r2'].iloc[0]
        return model_r2
    except (FileNotFoundError, IndexError, KeyError) as e:
        print(f"Error loading R2 for {best_model_name} from {METRICS_PATH}: {e}")
        return np.nan

def load_performance_summary(best_model_name):
    try:
        summary_df = pd.read_csv(PERF_SUMMARY_PATH)
        model_perf = summary_df[summary_df['model'] == best_model_name].iloc[0]
        # Add avg_return and volatility if not present (they are not by default from calculate_performance_metrics)
        # These will be calculated from the model's actual returns later
        if 'avg_return' not in model_perf:
            model_perf['avg_return'] = np.nan 
        if 'volatility' not in model_perf: # Using 'volatility' for monthly std
            model_perf['volatility'] = np.nan
        return model_perf
    except (FileNotFoundError, IndexError, KeyError) as e:
        print(f"Error loading performance summary for {best_model_name} from {PERF_SUMMARY_PATH}: {e}")
        return None

def get_spy_performance_for_period(model_returns_df):
    try:
        spy_returns_raw = pd.read_excel(SPY_DATA_PATH, index_col=0)
        spy_returns_raw.index = pd.to_datetime(spy_returns_raw.index)
        spy_returns_raw.columns = ['spy_ret'] # Ensure column name

        if model_returns_df is None or model_returns_df.empty:
            print("Model returns are empty, cannot align SPY data.")
            # Calculate SPY metrics over its entire history as a fallback
            spy_overall_perf = calculate_performance_metrics(spy_returns_raw, spy_returns_raw, 'SPY_overall') # Pass spy_returns_raw as benchmark too
            # We need to re-index it to match the expected output format for the table
            return pd.Series({
                'alpha': '—', # Alpha vs itself is 0, so put N/A
                'sharpe': f"{spy_overall_perf['sharpe'].iloc[0]:.2f}" if pd.notna(spy_overall_perf['sharpe'].iloc[0]) else 'N/A',
                'avg_return': f"{spy_returns_raw['spy_ret'].mean():.4f}",
                'volatility': f"{spy_returns_raw['spy_ret'].std():.4f}",
                'max_drawdown': f"{spy_overall_perf['max_drawdown'].iloc[0]:.4f}" if pd.notna(spy_overall_perf['max_drawdown'].iloc[0]) else 'N/A',
                'max_loss': f"{spy_overall_perf['max_loss'].iloc[0]:.4f}" if pd.notna(spy_overall_perf['max_loss'].iloc[0]) else 'N/A'
            })

        # Align SPY to model's OOS period
        model_start_date = model_returns_df.index.min()
        model_end_date = model_returns_df.index.max()
        spy_aligned = spy_returns_raw[(spy_returns_raw.index >= model_start_date) & (spy_returns_raw.index <= model_end_date)]
        
        if spy_aligned.empty:
            print("SPY returns have no overlap with model OOS period.")
            return None
        
        # Use calculate_performance_metrics for SPY against itself (to get drawdown etc.)
        # For Alpha, it will be 0. Avg Return and Volatility can be direct.
        spy_perf_metrics = calculate_performance_metrics(spy_aligned, spy_aligned, 'SPY') 

        return pd.Series({
            'alpha': '—', # Alpha of SPY vs SPY is not meaningful here, typically alpha is vs benchmark
            'sharpe': f"{spy_perf_metrics['sharpe'].iloc[0]:.2f}" if pd.notna(spy_perf_metrics['sharpe'].iloc[0]) else 'N/A',
            'avg_return': f"{spy_aligned['spy_ret'].mean():.4f}",
            'volatility': f"{spy_aligned['spy_ret'].std():.4f}",
            'max_drawdown': f"{spy_perf_metrics['max_drawdown'].iloc[0]:.4f}" if pd.notna(spy_perf_metrics['max_drawdown'].iloc[0]) else 'N/A',
            'max_loss': f"{spy_perf_metrics['max_loss'].iloc[0]:.4f}" if pd.notna(spy_perf_metrics['max_loss'].iloc[0]) else 'N/A'
        })
    except Exception as e:
        print(f"Error calculating SPY performance: {e}")
        return None

def update_presentation():
    best_model_name = get_best_model_name()
    if not best_model_name:
        return

    oos_r2 = load_model_metrics(best_model_name)
    model_perf_summary = load_performance_summary(best_model_name)

    model_port_ret_path = os.path.join(OUTPUT_DIR, f'{best_model_name}_overall_port_ret.csv')
    model_returns_df = None
    avg_model_monthly_return = np.nan
    std_model_monthly_vol = np.nan
    try:
        model_returns_df = pd.read_csv(model_port_ret_path, index_col=0, parse_dates=True)
        # Ensure the column is correctly identified (e.g. lasso_port_ret -> port_ret)
        model_returns_col = model_returns_df.columns[0]
        avg_model_monthly_return = model_returns_df[model_returns_col].mean()
        std_model_monthly_vol = model_returns_df[model_returns_col].std()
        if model_perf_summary is not None:
            model_perf_summary['avg_return'] = avg_model_monthly_return
            model_perf_summary['volatility'] = std_model_monthly_vol

    except FileNotFoundError:
        print(f"Model returns file not found: {model_port_ret_path}")
        
    spy_perf_for_table = get_spy_performance_for_period(model_returns_df)

    # --- Texts that might need your review based on final model --- 
    METHODOLOGY_TEMPLATE = (
        f"• Implemented {best_model_name.upper()} model to predict next-month stock excess returns\n"
        "• Training: Expanding window approach (initial 10-year: 2005-2014)\n"
        "• Validation: Expanding window (initial 2-year: 2015-2016)\n"
        "• OOS Testing: Expanding window (2017-2023), yearly re-estimation/re-training\n"
        "• Portfolio: Equal-weighted top 50 stocks monthly rebalancing\n"
        "• Features: 145 lagged stock characteristics (t) predicting returns (t+1)\n"
        "• Benchmark: S&P 500 (SPY) for performance comparison"
    )
    # You should customize this summary based on the chosen model and its specific results
    SUMMARY_TEMPLATE = (
        f"The {best_model_name.upper()} strategy demonstrated specific performance characteristics. "
        f"It achieved a monthly alpha of {model_perf_summary['alpha']:.4f} and an annualized Sharpe ratio of {model_perf_summary['sharpe']:.2f}. "
        f"The average monthly return was {model_perf_summary['avg_return']:.4f} and monthly volatility was {model_perf_summary['volatility']:.4f}, with a max drawdown of {model_perf_summary['max_drawdown']:.4f}. "
        f"The OOS R² was {oos_r2:.4f}. "
        "These results should be interpreted in context of the overall market performance and model limitations.\n"
        "Future improvements could involve further hyperparameter tuning, alternative feature sets, or different portfolio construction rules."
    )
    # --- End of texts needing review ---

    prs = Presentation(TEMPLATE_PATH)
    
    # Slide 1: Methodology
    slide = prs.slides[0]
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = METHODOLOGY_TEMPLATE
        p.font.size = Pt(18) # Adjusted font size
    
    # Slide 2: OOS R2
    slide = prs.slides[1]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = f"Out-of-sample $R^2$ ({best_model_name.upper()}): {oos_r2:.4f}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    r2_interp = "\nInterpretation:\n• R² measures the proportion of the variance for a dependent variable that's explained by an independent variable.\n"
    if pd.notna(oos_r2) and oos_r2 < 0:
        r2_interp += "• Negative R² indicates the model underperforms a simple mean prediction."
    elif pd.notna(oos_r2):
        r2_interp += f"• An R² of {oos_r2:.2f} suggests the model explains {oos_r2*100:.1f}% of the variance in returns."
    else:
        r2_interp += "• R² value not available."
    p.text = r2_interp
    p.font.size = Pt(18)
    
    # Slide 3: Performance Table
    slide = prs.slides[2]
    # Clear existing table first if any (assuming placeholder might contain one)
    for shape in slide.shapes:
        if shape.has_table:
            sp = shape
            slide.shapes._spTree.remove(sp._element)
            
    left = Inches(0.5); top = Inches(1.5); width = Inches(9.0); height = Inches(0.5)
    
    if model_perf_summary is not None and spy_perf_for_table is not None:
        perf_data = [
            ["Metric", f"{best_model_name.upper()} Strategy", "SPY"],
            ["Alpha (monthly)", f"{model_perf_summary['alpha']:.4f}", spy_perf_for_table['alpha']],
            ["Sharpe Ratio (ann.)", f"{model_perf_summary['sharpe']:.2f}", spy_perf_for_table['sharpe']],
            ["Avg Return (monthly)", f"{model_perf_summary['avg_return']:.4f}", spy_perf_for_table['avg_return']],
            ["Volatility (monthly)", f"{model_perf_summary['volatility']:.4f}", spy_perf_for_table['volatility']],
            ["Max Drawdown", f"{model_perf_summary['max_drawdown']:.4f}", spy_perf_for_table['max_drawdown']],
            ["Max 1-mo Loss", f"{model_perf_summary['max_loss']:.4f}", spy_perf_for_table['max_loss']]
        ]
        rows, cols = len(perf_data), len(perf_data[0])
        table = slide.shapes.add_table(rows, cols, left, top, width, Inches(rows * 0.4)).table
        for r, row_data in enumerate(perf_data):
            for c, cell_data in enumerate(row_data):
                cell = table.cell(r,c)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                if r == 0: cell.text_frame.paragraphs[0].font.bold = True
    else:
        tf = slide.shapes.placeholders[1].text_frame; tf.clear(); tf.add_paragraph().text = "Performance data not available."

    # Slide 4: Cumulative Returns Plot
    slide = prs.slides[3]
    if os.path.exists(CUM_RET_PNG):
        # Clear existing pictures first (if any)
        for shape in slide.shapes:
            if shape.shape_type == 13: # msoPicture
                 sp = shape
                 slide.shapes._spTree.remove(sp._element)
        slide.shapes.add_picture(CUM_RET_PNG, Inches(0.5), Inches(1.0), width=Inches(9.0))
    else:
        tf = slide.shapes.placeholders[1].text_frame; tf.clear(); tf.add_paragraph().text = "Cumulative returns plot not found."
    
    # Slide 5: Summary
    slide = prs.slides[4]
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = SUMMARY_TEMPLATE if model_perf_summary is not None else "Summary data not available."
        p.font.size = Pt(16) # Adjusted font size
    
    output_ppt_path = os.path.join(OUTPUT_DIR, f'Module3_Report_{best_model_name}.pptx')
    prs.save(output_ppt_path)
    print(f"Presentation updated and saved to: {output_ppt_path}")

if __name__ == "__main__":
    update_presentation() 